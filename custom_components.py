from haystack import component, Document
from typing import List, Union, Dict, Optional, Any
from pathlib import Path
import json
from boto3 import Session
from haystack.dataclasses import ChatMessage, ChatRole
from typing import ClassVar
from abc import ABC, abstractmethod
import docx
import pandas as pd
from pptx import Presentation
import boto3
import os
from copy import deepcopy
from typing import Dict, List, Literal, Tuple

from more_itertools import windowed


AI_CENTRAL_ENABLED = os.getenv("ai_central_enabled", "false")
AI_CENTRAL_LAMBDA_FUNCTION_ARN = os.getenv("ai_central_lambda_function_arn")

LAMBDA_CLIENT = boto3.client("lambda")
BEDROCK_CLIENT = boto3.client("bedrock-runtime")

def get_response_body_from_ai_central_response(ai_central_response):
    response = ai_central_response.get("Payload").read().decode()
    return json.loads(json.loads(json.loads(response).get("body")))

def normalize_metadata(
    meta: Optional[Union[Dict[str, Any], List[Dict[str, Any]]]], sources_count: int
) -> List[Dict[str, Any]]:
        if meta is None:
            return [{}] * sources_count
        if isinstance(meta, dict):
            return [meta] * sources_count
        if isinstance(meta, list):
            if sources_count != len(meta):
                raise ValueError("The length of the metadata list must match the number of sources.")
            return meta
        raise ValueError("meta must be either None, a dictionary or a list of dictionaries.")

@component
class DocxToTextConverter:
    """
    A component generating personal welcome message and making it upper case
    """

    @component.output_types(documents=List[Document])
    def run(
        self,
        sources: List[Union[str, Path]],
        meta: Optional[Union[Dict[str, Any], List[Dict[str, Any]]]] = None,
    ):
        documents = []
        meta_list = normalize_metadata(meta=meta, sources_count=len(sources))
        for source, metadata in zip(sources, meta_list):
            file = docx.Document(source)
            paragraphs = [para.text for para in file.paragraphs]
            text = "\n".join(paragraphs)
            document = Document(content=text, meta=metadata)
            documents.append(document)
        return {"documents": documents}


@component
class ExcelToMarkdown:
    @component.output_types(documents=List[Document])
    def run(
        self,
        sources: List[Union[str, Path]],
        meta: Optional[Union[Dict[str, Any], List[Dict[str, Any]]]] = None,
    ):
        documents = []
        meta_list = normalize_metadata(meta=meta, sources_count=len(sources))
        for source, metadata in zip(sources, meta_list):
            if source.split(".")[-1] == "xlsx":
                df = pd.read_excel(source, sheet_name=None)
                for key in df.keys():
                    for i in range(0, df[key].shape[0], 5):
                        end_row = min(i + 5, df[key].shape[0])
                        markdown = (df[key].iloc[i:end_row, :]).to_markdown()
                        document = Document(content=markdown, meta=metadata)
                        documents.append(document)
            else:
                df = pd.read_csv(source)

                for i in range(0, df.shape[0], 5):
                    end_row = min(i + 5, df.shape[0])
                    markdown = (df.iloc[i:end_row, :]).to_markdown()

                    document = Document(content=markdown, meta=metadata)
                    documents.append(document)

        return {"documents": documents}


@component
class PptxConverter:
    @component.output_types(documents=List[Document])
    def run(
        self,
        sources: List[Union[str, Path]],
        meta: Optional[Union[Dict[str, Any], List[Dict[str, Any]]]] = None,
    ):
        documents = []
        meta_list = normalize_metadata(meta=meta, sources_count=len(sources))
        for source, metadata in zip(sources, meta_list):
            pres = Presentation(source)
            text_parts = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            text = "\n".join(text_parts)
            documents.append(Document(content=text, meta=metadata))

        return {"documents": documents}
